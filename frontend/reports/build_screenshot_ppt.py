from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def main() -> None:
    images_dir = Path("/Users/rushan/Library/Application Support/Code/User/workspaceStorage/vscode-chat-images")
    out_file = Path("/Users/rushan/Main/MERN/QURE/frontend/reports/QURE-presentation-from-screenshots.pptx")

    # Keep only the screenshots from the latest deck shared in chat.
    names = [
        "image-1775043731166.png",
        "image-1775043741561.png",
        "image-1775043749194.png",
        "image-1775043759255.png",
        "image-1775043767109.png",
        "image-1775043773800.png",
        "image-1775043782981.png",
        "image-1775043789816.png",
        "image-1775043797213.png",
        "image-1775043804328.png",
    ]

    image_paths = [images_dir / name for name in names]
    missing = [str(path) for path in image_paths if not path.exists()]
    if missing:
        raise FileNotFoundError("Missing screenshot files:\n" + "\n".join(missing))

    prs = Presentation()

    # Match screenshot aspect ratio for exact visual fidelity.
    # 2940 x 1912 -> ratio ~1.53766
    slide_width = Inches(13.333)
    slide_height = Inches(13.333 * 1912 / 2940)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    blank_layout = prs.slide_layouts[6]

    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    out_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_file))
    print(f"Created: {out_file}")


if __name__ == "__main__":
    main()
